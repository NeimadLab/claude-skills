"""
PISA Programmatic QA Engine
Reads a generated PPTX back and validates it against PISA rules.
Zero external dependencies beyond python-pptx.
Replaces all LibreOffice-based visual QA.

Usage:
    from programmatic_qa import qa_deck, qa_slide, qa_report
    issues = qa_deck("output.pptx", theme, persona=None)
    print(qa_report(issues))
"""

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re, json

# ═══════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════

SLIDE_W_EMU = 12192000  # 16:9
SLIDE_H_EMU = 6858000
TOKEN_PATTERN = re.compile(r'token\.\w+\.\w+')
SEVERITY = {"critical": 3, "warning": 2, "info": 1}

# Default density limits (overridden by persona)
DEFAULT_LIMITS = {
    "max_words_body": 75,
    "max_bullets": 5,
    "max_hierarchy_levels": 2,
    "max_content_shapes": 6,
    "max_total_shapes": 8,
}


# ═══════════════════════════════════════════════════
# INDIVIDUAL CHECKS
# ═══════════════════════════════════════════════════

def check_token_leaks(slide, slide_num):
    """Find any text still containing unresolved token references."""
    issues = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                matches = TOKEN_PATTERN.findall(run.text)
                if matches:
                    issues.append({
                        "slide": slide_num,
                        "check": "token_leak",
                        "severity": "critical",
                        "shape": shape.name,
                        "detail": f"Unresolved tokens: {', '.join(matches)}",
                        "fix": "Theme JSON is missing these keys. Add them and re-render."
                    })
    # Also check shape fills for token strings (shouldn't happen but defensive)
    for shape in slide.shapes:
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                fc = str(shape.fill.fore_color.rgb) if shape.fill.fore_color else ""
                if "token" in fc.lower():
                    issues.append({
                        "slide": slide_num, "check": "token_leak", "severity": "critical",
                        "shape": shape.name, "detail": f"Fill colour is a token ref: {fc}",
                        "fix": "Token resolution failed for this shape's fill."
                    })
        except Exception:
            pass
    return issues


def check_overlaps(slide, slide_num):
    """Detect overlapping content shapes (excluding decorative thin elements)."""
    issues = []
    shapes = []
    for s in slide.shapes:
        if s.width < 91440 or s.height < 91440:  # skip thin decorative
            continue
        shapes.append(s)

    for i, a in enumerate(shapes):
        for j, b in enumerate(shapes):
            if j <= i:
                continue
            # Compute overlap
            ax1, ay1 = a.left, a.top
            ax2, ay2 = a.left + a.width, a.top + a.height
            bx1, by1 = b.left, b.top
            bx2, by2 = b.left + b.width, b.top + b.height

            ox = max(0, min(ax2, bx2) - max(ax1, bx1))
            oy = max(0, min(ay2, by2) - max(ay1, by1))
            overlap_area = ox * oy

            if overlap_area > 0:
                # Calculate overlap as percentage of smaller shape
                area_a = a.width * a.height
                area_b = b.width * b.height
                smaller = min(area_a, area_b)
                overlap_pct = (overlap_area / smaller * 100) if smaller > 0 else 0

                if overlap_pct > 15:
                    issues.append({
                        "slide": slide_num, "check": "overlap", "severity": "warning",
                        "shape": f"{a.name} × {b.name}",
                        "detail": f"Shapes overlap by {overlap_pct:.0f}% of smaller shape",
                        "fix": "Reposition shapes to eliminate overlap, or merge into one component."
                    })
    return issues


def check_out_of_bounds(slide, slide_num):
    """Detect shapes that extend beyond slide boundaries."""
    issues = []
    margin = 91440  # 0.1" tolerance
    for s in slide.shapes:
        problems = []
        if s.left < -margin:
            problems.append("extends left of slide")
        if s.top < -margin:
            problems.append("extends above slide")
        if s.left + s.width > SLIDE_W_EMU + margin:
            problems.append("extends right of slide")
        if s.top + s.height > SLIDE_H_EMU + margin:
            problems.append("extends below slide")
        if problems:
            issues.append({
                "slide": slide_num, "check": "out_of_bounds", "severity": "warning",
                "shape": s.name, "detail": "; ".join(problems),
                "fix": "Adjust shape position to stay within slide boundaries."
            })
    return issues


def check_empty_shapes(slide, slide_num):
    """Detect shapes that should have content but are empty."""
    issues = []
    for s in slide.shapes:
        if s.width < 274320 and s.height < 274320:  # skip tiny decorative
            continue
        if hasattr(s, 'text_frame'):
            text = s.text_frame.text.strip()
            # Large text shape with no content
            if not text and s.width > 914400 and s.height > 457200:
                issues.append({
                    "slide": slide_num, "check": "empty_shape", "severity": "info",
                    "shape": s.name,
                    "detail": f"Large text shape ({s.width/914400:.1f}\" × {s.height/914400:.1f}\") is empty",
                    "fix": "Add content or remove the placeholder shape."
                })
    return issues


def check_text_overflow(slide, slide_num):
    """Estimate if text is likely to overflow its container."""
    issues = []
    for s in slide.shapes:
        if not hasattr(s, 'text_frame'):
            continue
        text = s.text_frame.text.strip()
        if not text:
            continue

        # Estimate character capacity based on shape width
        # Approximate: 7pt per character at 12pt font, with margins
        shape_w_inches = s.width / 914400
        shape_h_inches = s.height / 914400

        # Get font size (default 12pt if not set)
        max_font_pt = 12
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    fs = run.font.size.pt
                    if fs > max_font_pt:
                        max_font_pt = fs

        chars_per_line = int((shape_w_inches - 0.2) / (max_font_pt * 0.0075))
        if chars_per_line <= 0:
            chars_per_line = 1
        lines_available = int((shape_h_inches - 0.15) / (max_font_pt * 0.018))
        if lines_available <= 0:
            lines_available = 1
        max_chars = chars_per_line * lines_available

        if len(text) > max_chars * 1.3:  # 30% tolerance
            issues.append({
                "slide": slide_num, "check": "text_overflow", "severity": "warning",
                "shape": s.name,
                "detail": f"Text ({len(text)} chars) likely overflows shape (est. capacity: {max_chars} chars)",
                "fix": "Reduce text length, decrease font size, or increase shape dimensions."
            })
    return issues


def check_density(slide, slide_num, limits=None):
    """Check slide density against limits (from persona or defaults)."""
    if limits is None:
        limits = DEFAULT_LIMITS

    issues = []
    total_words = 0
    content_shapes = 0
    total_shapes = 0
    max_bullets_found = 0

    for s in slide.shapes:
        total_shapes += 1
        if s.width < 91440 or s.height < 91440:
            continue
        content_shapes += 1

        if hasattr(s, 'text_frame'):
            text = s.text_frame.text.strip()
            total_words += len(text.split())
            # Count bullets (paragraphs in the same text frame)
            bullet_count = sum(1 for p in s.text_frame.paragraphs if p.text.strip())
            if bullet_count > max_bullets_found:
                max_bullets_found = bullet_count

    if total_words > limits.get("max_words_body", 75):
        issues.append({
            "slide": slide_num, "check": "density_words", "severity": "warning",
            "shape": "slide",
            "detail": f"Total words: {total_words} (limit: {limits['max_words_body']})",
            "fix": "Split into two slides or move detail to speaker notes."
        })

    if max_bullets_found > limits.get("max_bullets", 5):
        issues.append({
            "slide": slide_num, "check": "density_bullets", "severity": "warning",
            "shape": "slide",
            "detail": f"Max bullets in one group: {max_bullets_found} (limit: {limits['max_bullets']})",
            "fix": "Consolidate items or restructure into a different layout."
        })

    if content_shapes > limits.get("max_content_shapes", 6):
        issues.append({
            "slide": slide_num, "check": "density_shapes", "severity": "info",
            "shape": "slide",
            "detail": f"Content shapes: {content_shapes} (limit: {limits['max_content_shapes']})",
            "fix": "Simplify or split the slide."
        })

    return issues


def check_title_quality(slide, slide_num):
    """Evaluate the title shape against the 5 title rules."""
    issues = []
    title_shape = None

    # Find the title (topmost wide text shape)
    for s in slide.shapes:
        if not hasattr(s, 'text_frame'):
            continue
        if s.top < SLIDE_H_EMU * 0.2 and s.width > SLIDE_W_EMU * 0.4:
            text = s.text_frame.text.strip()
            if text:
                title_shape = s
                break

    if title_shape is None:
        issues.append({
            "slide": slide_num, "check": "title_missing", "severity": "warning",
            "shape": "slide", "detail": "No title shape detected",
            "fix": "Add a title — it anchors the slide's message."
        })
        return issues

    title_text = title_shape.text_frame.text.strip()
    words = title_text.split()

    # Rule 1: Insight, not label (heuristic: ≤3 words with no verb = likely label)
    label_patterns = [
        r'^(overview|summary|introduction|background|agenda|appendix|conclusion)$',
        r'^[\w\s]{1,20}(overview|update|status|report|analysis|review)$',
    ]
    is_label = len(words) <= 3 and not any(c in title_text for c in [' is ', ' are ', ' was ', ' grew ', ' increased ', ' decreased ', ' shows '])
    for pat in label_patterns:
        if re.match(pat, title_text, re.IGNORECASE):
            is_label = True
            break
    if is_label and len(words) <= 3:
        issues.append({
            "slide": slide_num, "check": "title_is_label", "severity": "info",
            "shape": title_shape.name,
            "detail": f"Title \"{title_text}\" appears to be a label, not an insight",
            "fix": "Rewrite as an insight sentence: what is the takeaway?"
        })

    # Rule 2: Under 12 words
    if len(words) > 12:
        issues.append({
            "slide": slide_num, "check": "title_too_long", "severity": "info",
            "shape": title_shape.name,
            "detail": f"Title is {len(words)} words (recommended: ≤12)",
            "fix": "Shorten to a concise insight statement."
        })

    # Rule 3: No ending punctuation
    if title_text and title_text[-1] in '.!?':
        issues.append({
            "slide": slide_num, "check": "title_punctuation", "severity": "info",
            "shape": title_shape.name,
            "detail": f"Title ends with '{title_text[-1]}' — titles should not have ending punctuation",
            "fix": "Remove the trailing punctuation."
        })

    return issues


def check_colour_consistency(slide, slide_num, theme=None):
    """Verify that colours used match the expected theme palette."""
    if theme is None:
        return []

    issues = []
    theme_colors = set()
    for k, v in theme.items():
        if k.startswith("token.color."):
            theme_colors.add(v.upper().lstrip("#"))

    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb:
                            c = str(run.font.color.rgb).upper()
                            if c not in theme_colors and c not in ("000000", "FFFFFF", ""):
                                issues.append({
                                    "slide": slide_num, "check": "off_palette_colour",
                                    "severity": "info", "shape": s.name,
                                    "detail": f"Text colour #{c} is not in the theme palette",
                                    "fix": "Use a theme token colour instead."
                                })
                                break  # one per shape is enough
                    except Exception:
                        pass
    return issues


# ═══════════════════════════════════════════════════
# V2.1 TEMPLATE QA (validates the template JSON before rendering)
# ═══════════════════════════════════════════════════

def qa_template(prim, slide_num=0):
    """Validate a V2.1 template JSON for completeness. Returns list of issues."""
    issues = []
    
    # Check schema version
    sv = prim.get("schema_version", 2)
    
    for i, comp in enumerate(prim.get("components", [])):
        role = comp.get("role", "")
        
        # V2.1: KPI cards should have kpi{}
        if role == "kpi_card" and not comp.get("kpi"):
            issues.append({
                "slide": slide_num, "check": "missing_kpi_data", "severity": "warning",
                "shape": f"component[{i}]",
                "detail": f"KPI card has no kpi{{}} — value/unit/label will be guessed from text_preview",
                "fix": "Add kpi: {{value, unit, label}} to this component."
            })
        
        # V2.1: All components should have visual{}
        if not comp.get("visual") and role not in ("footer", "decoration"):
            issues.append({
                "slide": slide_num, "check": "missing_visual", "severity": "info",
                "shape": f"component[{i}] ({role})",
                "detail": "No visual{{}} — bg_variant and icon_hint will default to generic",
                "fix": "Add visual: {{bg_variant, has_icon, icon_hint}} for richer rendering."
            })
        
        # V2.1: Body with timeline/org/checklist content should have content_data{}
        if role == "body" and not comp.get("content_data"):
            tp = comp.get("text_preview", "")
            # Heuristic: if text has structured patterns, it probably needs content_data
            if any(p in tp for p in ["1.", "2.", "3.", "→", "✓", "•"]):
                issues.append({
                    "slide": slide_num, "check": "flat_structured_content", "severity": "info",
                    "shape": f"component[{i}]",
                    "detail": "Body text appears to have list/structured content but no content_data{{}}",
                    "fix": "Add content_data with type and entries[] for proper rendering."
                })
    
    # V2.1: Compound slides should declare zones
    si = prim.get("secondary_intent")
    if si and not prim.get("zones"):
        issues.append({
            "slide": slide_num, "check": "compound_no_zones", "severity": "warning",
            "shape": "template",
            "detail": f"Has secondary_intent '{si}' but no zones[] to map components to intents",
            "fix": "Add zones[] with zone_intent and component_indices for each region."
        })
    
    # Check layout_variant exists
    if not prim.get("layout_variant"):
        issues.append({
            "slide": slide_num, "check": "missing_layout_variant", "severity": "info",
            "shape": "template",
            "detail": "No layout_variant — generic rendering will be used",
            "fix": "Add layout_variant (e.g. 'grid_alternating_dark_light') for pattern-specific rendering."
        })
    
    return issues


# ═══════════════════════════════════════════════════
# ORCHESTRATORS
# ═══════════════════════════════════════════════════

def qa_slide(slide, slide_num, theme=None, persona_limits=None):
    """Run all QA checks on a single slide. Returns list of issue dicts."""
    issues = []
    issues.extend(check_token_leaks(slide, slide_num))
    issues.extend(check_overlaps(slide, slide_num))
    issues.extend(check_out_of_bounds(slide, slide_num))
    issues.extend(check_empty_shapes(slide, slide_num))
    issues.extend(check_text_overflow(slide, slide_num))
    issues.extend(check_density(slide, slide_num, persona_limits))
    issues.extend(check_title_quality(slide, slide_num))
    issues.extend(check_colour_consistency(slide, slide_num, theme))
    return issues


def qa_deck(pptx_path, theme=None, persona_limits=None):
    """Run all QA checks on an entire PPTX deck.

    Args:
        pptx_path: path to the PPTX file
        theme: theme JSON dict (for colour consistency check)
        persona_limits: density limit overrides from persona

    Returns:
        dict with per-slide issues and deck-level summary
    """
    prs = Presentation(pptx_path)
    all_issues = []
    slide_summaries = []

    for i, slide in enumerate(prs.slides):
        slide_issues = qa_slide(slide, i + 1, theme, persona_limits)
        all_issues.extend(slide_issues)
        slide_summaries.append({
            "slide": i + 1,
            "issues": len(slide_issues),
            "critical": sum(1 for x in slide_issues if x["severity"] == "critical"),
            "warnings": sum(1 for x in slide_issues if x["severity"] == "warning"),
            "info": sum(1 for x in slide_issues if x["severity"] == "info"),
        })

    # Deck-level checks
    deck_issues = []

    # Check for duplicate titles
    titles = []
    for slide in prs.slides:
        for s in slide.shapes:
            if hasattr(s, 'text_frame') and s.top < SLIDE_H_EMU * 0.2 and s.width > SLIDE_W_EMU * 0.4:
                t = s.text_frame.text.strip()
                if t:
                    titles.append(t)
                    break

    seen = {}
    for t in titles:
        key = t.lower().strip()
        if key in seen:
            deck_issues.append({
                "slide": 0, "check": "duplicate_title", "severity": "warning",
                "shape": "deck", "detail": f"Title \"{t}\" appears multiple times",
                "fix": "Each slide should have a unique title."
            })
        seen[key] = True

    # Check slide count sanity
    n = len(prs.slides)
    if n > 50:
        deck_issues.append({
            "slide": 0, "check": "deck_too_long", "severity": "info",
            "shape": "deck", "detail": f"Deck has {n} slides — consider splitting or adding appendix markers",
            "fix": "Move supporting material to appendix sections."
        })

    all_issues.extend(deck_issues)

    return {
        "total_slides": len(prs.slides),
        "total_issues": len(all_issues),
        "critical": sum(1 for x in all_issues if x["severity"] == "critical"),
        "warnings": sum(1 for x in all_issues if x["severity"] == "warning"),
        "info": sum(1 for x in all_issues if x["severity"] == "info"),
        "pass": sum(1 for x in all_issues if x["severity"] == "critical") == 0,
        "slide_summaries": slide_summaries,
        "issues": all_issues,
    }


def qa_report(result):
    """Format QA results as a readable text report."""
    lines = []
    lines.append("╔══════════════════════════════════════════════════════╗")
    status = "✅ PASS" if result["pass"] else "❌ FAIL"
    lines.append(f"║  PISA QA Report — {status}                        ║")
    lines.append("╠══════════════════════════════════════════════════════╣")
    lines.append(f"║  Slides: {result['total_slides']}    Issues: {result['total_issues']}                        ║")
    lines.append(f"║  Critical: {result['critical']}   Warnings: {result['warnings']}   Info: {result['info']}          ║")
    lines.append("╚══════════════════════════════════════════════════════╝")
    lines.append("")

    for s in result["slide_summaries"]:
        if s["issues"] > 0:
            flag = "🔴" if s["critical"] > 0 else ("🟡" if s["warnings"] > 0 else "🔵")
            lines.append(f"  {flag} Slide {s['slide']}: {s['issues']} issue(s)")

    lines.append("")
    for issue in result["issues"]:
        sev = {"critical": "🔴", "warning": "🟡", "info": "🔵"}[issue["severity"]]
        slide = f"Slide {issue['slide']}" if issue["slide"] > 0 else "Deck"
        lines.append(f"  {sev} [{issue['check']}] {slide} — {issue['shape']}")
        lines.append(f"     {issue['detail']}")
        lines.append(f"     Fix: {issue['fix']}")
        lines.append("")

    return "\n".join(lines)


# ═══════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="PISA Programmatic QA")
    parser.add_argument("pptx", help="Path to PPTX file")
    parser.add_argument("--theme", help="Path to theme JSON")
    parser.add_argument("--persona-limits", help="JSON string of density limits")
    args = parser.parse_args()

    theme = None
    if args.theme:
        with open(args.theme) as f:
            theme = json.load(f)

    limits = None
    if args.persona_limits:
        limits = json.loads(args.persona_limits)

    result = qa_deck(args.pptx, theme, limits)
    print(qa_report(result))
