"""
PISA V2 — Slide-to-Primitive Extraction Engine
Handles non-grid, freeform, rotated, grouped, and overlapping shapes.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import math, json, hashlib, datetime, os, re
import numpy as np
from scipy.spatial import KDTree
from sklearn.cluster import DBSCAN

# ═══════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════
SLIDE_W = 12192000   # 16:9 widescreen EMU
SLIDE_H = 6858000
DECORATIVE_MIN = 91440   # 0.1" — thinner = decorative
TINY_SHAPE = 274320      # 0.3" — smaller on both axes = icon/bullet
EDGE_TOLERANCE = 1.5     # % for alignment detection
DEDUP_THRESHOLD = 0.82
DBSCAN_EPS = 6.0         # reading-order band detection
KPI_PROXIMITY = 15.0     # % distance for KPI card grouping


# ═══════════════════════════════════════════════════
# STEP 1 — DEEP SHAPE INVENTORY
# ═══════════════════════════════════════════════════

def _oriented_bbox(left, top, width, height, rotation_deg):
    if rotation_deg == 0:
        return left, top, width, height
    cx, cy = left + width/2, top + height/2
    rad = math.radians(rotation_deg)
    cos_a, sin_a = abs(math.cos(rad)), abs(math.sin(rad))
    nw = width * cos_a + height * sin_a
    nh = width * sin_a + height * cos_a
    return int(cx - nw/2), int(cy - nh/2), int(nw), int(nh)

def _flatten_group(group_shape, px=0, py=0):
    results = []
    for child in group_shape.shapes:
        ax, ay = px + child.left, py + child.top
        if child.shape_type == MSO_SHAPE_TYPE.GROUP:
            results.extend(_flatten_group(child, ax, ay))
        else:
            results.append({"shape": child, "abs_left": ax, "abs_top": ay,
                            "width": child.width, "height": child.height, "from_group": True})
    return results

def _is_decorative(shape, w, h):
    has_text = getattr(shape, 'has_text_frame', False) and bool(
        getattr(shape.text_frame, 'text', '').strip() if hasattr(shape, 'text_frame') else '')
    if (w < DECORATIVE_MIN or h < DECORATIVE_MIN) and not has_text:
        return True
    if w < TINY_SHAPE and h < TINY_SHAPE and not has_text:
        return True
    return False

def deep_inventory(slide, slide_index):
    raw = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            raw.extend(_flatten_group(shape, shape.left, shape.top))
        else:
            raw.append({"shape": shape, "abs_left": shape.left, "abs_top": shape.top,
                        "width": shape.width, "height": shape.height, "from_group": False})

    inventory = []
    for item in raw:
        s = item["shape"]
        rot = (getattr(s, 'rotation', 0) or 0) % 360
        bl, bt, bw, bh = _oriented_bbox(item["abs_left"], item["abs_top"],
                                          item["width"], item["height"], rot)
        entry = {
            "shape_id": s.shape_id, "name": s.name,
            "left": bl, "top": bt, "width": bw, "height": bh,
            "rotation": rot, "from_group": item["from_group"],
            "decorative": _is_decorative(s, bw, bh),
            "has_text": False, "has_table": getattr(s, 'has_table', False),
            "has_chart": hasattr(s, 'chart'),
            "has_image": (s.shape_type == MSO_SHAPE_TYPE.PICTURE
                          if hasattr(s, 'shape_type') else False),
            "x_pct": round(bl / SLIDE_W * 100, 2),
            "y_pct": round(bt / SLIDE_H * 100, 2),
            "w_pct": round(bw / SLIDE_W * 100, 2),
            "h_pct": round(bh / SLIDE_H * 100, 2),
        }
        if getattr(s, 'has_text_frame', False):
            text = s.text_frame.text.strip()
            if text:
                entry["has_text"] = True
                entry["text"] = text[:200]
                entry["word_count"] = len(text.split())
        inventory.append(entry)
    return inventory


# ═══════════════════════════════════════════════════
# STEP 2 — SPATIAL ANALYSIS & LAYOUT CLASSIFICATION
# ═══════════════════════════════════════════════════

def _find_axis_clusters(values, tol=EDGE_TOLERANCE):
    if len(values) < 2:
        return []
    sv = sorted(enumerate(values), key=lambda x: x[1])
    clusters, cur = [], [sv[0]]
    for i in range(1, len(sv)):
        if sv[i][1] - cur[-1][1] <= tol:
            cur.append(sv[i])
        else:
            if len(cur) >= 2:
                clusters.append({"value": np.mean([v for _, v in cur]),
                                  "indices": [idx for idx, _ in cur]})
            cur = [sv[i]]
    if len(cur) >= 2:
        clusters.append({"value": np.mean([v for _, v in cur]),
                          "indices": [idx for idx, _ in cur]})
    return clusters

def detect_alignment_axes(shapes):
    lefts   = [s["x_pct"] for s in shapes]
    rights  = [s["x_pct"] + s["w_pct"] for s in shapes]
    tops    = [s["y_pct"] for s in shapes]
    bottoms = [s["y_pct"] + s["h_pct"] for s in shapes]
    cx      = [s["x_pct"] + s["w_pct"]/2 for s in shapes]
    cy      = [s["y_pct"] + s["h_pct"]/2 for s in shapes]
    return {
        "left": _find_axis_clusters(lefts), "right": _find_axis_clusters(rights),
        "top": _find_axis_clusters(tops), "bottom": _find_axis_clusters(bottoms),
        "center_x": _find_axis_clusters(cx), "center_y": _find_axis_clusters(cy),
    }

def classify_layout(shapes, axes):
    n = len(shapes)
    if n <= 1:
        return "single"
    has_col = any(len(a["indices"]) >= 2 for a in axes["top"])
    has_row = any(len(a["indices"]) >= 2 for a in axes["left"])
    x_al = len(axes["left"]) + len(axes["center_x"])
    y_al = len(axes["top"]) + len(axes["center_y"])

    if has_col and has_row and x_al >= 2 and y_al >= 2:
        return "grid"
    if has_col and x_al < y_al:
        return "columns"
    if has_row and y_al < x_al:
        return "rows"

    cx = [s["x_pct"] + s["w_pct"]/2 for s in shapes]
    cy = [s["y_pct"] + s["h_pct"]/2 for s in shapes]
    if (max(cx)-min(cx)) < 30 and (max(cy)-min(cy)) > 30:
        return "stacked"

    mcx, mcy = np.mean(cx), np.mean(cy)
    dists = [((x-mcx)**2 + (y-mcy)**2)**0.5 for x, y in zip(cx, cy)]
    if n >= 3 and np.std(dists) / (np.mean(dists) + 0.01) < 0.3:
        return "radial"

    return "freeform"

def infer_reading_order(shapes):
    if not shapes:
        return []
    yc = np.array([[s["y_pct"] + s["h_pct"]/2] for s in shapes])
    labels = DBSCAN(eps=DBSCAN_EPS, min_samples=1).fit(yc).labels_
    bands = {}
    for i, l in enumerate(labels):
        bands.setdefault(l, []).append(i)
    ordered = []
    for l in sorted(bands, key=lambda lb: np.mean([yc[i][0] for i in bands[lb]])):
        ordered.extend(sorted(bands[l], key=lambda i: shapes[i]["x_pct"]))
    return ordered


# ═══════════════════════════════════════════════════
# STEP 3 — ROLE ASSIGNMENT
# ═══════════════════════════════════════════════════

def assign_roles(shapes, reading_order, layout_type):
    roles = [None] * len(shapes)
    for idx in reading_order:
        s = shapes[idx]
        if s.get("decorative"):  roles[idx] = "decoration"; continue
        if s.get("has_chart"):   roles[idx] = "chart"; continue
        if s.get("has_table"):   roles[idx] = "table"; continue
        if s.get("has_image"):   roles[idx] = "image"; continue
        if not s.get("has_text"):roles[idx] = "shape"; continue

        wc, y, h, w = s.get("word_count", 0), s["y_pct"], s["h_pct"], s["w_pct"]
        if y < 18 and w > 40 and wc <= 15 and h < 20 and "title" not in roles:
            roles[idx] = "title"; continue
        if y < 25 and wc <= 20 and "title" in roles and roles[idx] is None:
            roles[idx] = "subtitle"; continue
        if y > 85 and wc <= 15:
            roles[idx] = "footer"; continue
        if wc <= 6 and w < 35 and h < 30 and re.search(r'\d', s.get("text", "")):
            roles[idx] = "kpi_card"; continue
        if wc <= 4 and (w < 25 or h < 12):
            roles[idx] = "label"; continue
        roles[idx] = "body"
    return roles


# ═══════════════════════════════════════════════════
# STEP 4 — SEMANTIC GROUP DETECTION
# ═══════════════════════════════════════════════════

def detect_semantic_groups(shapes, roles, reading_order):
    groups, used = [], set()
    content = [(i, shapes[i]) for i in reading_order if not shapes[i].get("decorative")]

    for _, (idx, s) in enumerate(content):
        if idx in used:
            continue
        if roles[idx] == "kpi_card":
            g = {"type": "kpi_unit", "members": [idx]}
            cx, cy = s["x_pct"]+s["w_pct"]/2, s["y_pct"]+s["h_pct"]/2
            for _, (jdx, t) in enumerate(content):
                if jdx == idx or jdx in used:
                    continue
                d = ((cx-(t["x_pct"]+t["w_pct"]/2))**2 + (cy-(t["y_pct"]+t["h_pct"]/2))**2)**0.5
                if d < KPI_PROXIMITY and roles[jdx] in ("label", "shape", "image"):
                    g["members"].append(jdx)
            if len(g["members"]) > 1:
                used.update(g["members"]); groups.append(g); continue
        if roles[idx] == "image":
            g = {"type": "image_caption", "members": [idx]}
            for _, (jdx, t) in enumerate(content):
                if jdx == idx or jdx in used:
                    continue
                if roles[jdx] == "label" and abs(t["x_pct"] - s["x_pct"]) < 15:
                    if abs((t["y_pct"]+t["h_pct"]/2) - (s["y_pct"]+s["h_pct"]/2)) < 20:
                        g["members"].append(jdx)
            if len(g["members"]) > 1:
                used.update(g["members"]); groups.append(g)

    for idx in reading_order:
        if idx not in used and not shapes[idx].get("decorative"):
            groups.append({"type": "standalone", "members": [idx]}); used.add(idx)
    return groups


# ═══════════════════════════════════════════════════
# STEP 5 — INTENT CLASSIFICATION (TWO-PASS)
# ═══════════════════════════════════════════════════

INTENT_RULES = {
    "cover":              {"max_content": 4, "max_words": 25, "first": True},
    "executive_summary":  {"keywords": ["summary","overview","highlights","key"], "word_range": (20,80)},
    "kpi_dashboard":      {"min_kpi": 2, "keywords": ["KPI","%","YoY","vs","revenue","margin"]},
    "comparison_columns": {"layouts": ["columns"], "keywords": ["vs","compare","option"]},
    "linear_process":     {"layouts": ["rows","columns"], "keywords": ["step","phase","stage"], "min_n": 3},
    "timeline":           {"keywords": ["Q1","Q2","Q3","Q4","2024","2025","2026","milestone","roadmap"]},
    "matrix_2x2":         {"layouts": ["grid"], "n_exact": 4, "keywords": ["quadrant","matrix"]},
    "single_insight":     {"max_content": 3, "max_words": 45},
    "section_divider":    {"max_content": 2, "max_words": 15},
    "conclusion_cta":     {"keywords": ["next steps","action","recommendation","conclusion"]},
    "data_table":         {"has_table": True},
    "chart_driven":       {"has_chart": True},
    "agenda":             {"keywords": ["agenda","contents","outline","topics"]},
    "thank_you":          {"keywords": ["thank","questions","Q&A","contact"], "last": True},
    "quote_testimonial":  {"keywords": ["quote","said"], "max_content": 3},
    "before_after":       {"keywords": ["before","after","current","future","as-is","to-be"]},
    "swot":               {"keywords": ["SWOT","strength","weakness","opportunity","threat"], "layouts": ["grid"]},
    "funnel":             {"keywords": ["funnel","pipeline","conversion"]},
    "org_chart":          {"min_n": 5, "keywords": ["team","organization","reporting"]},
    "image_showcase":     {"dom_image": True, "max_words": 25},
}

def classify_intent_v2(shapes, roles, groups, layout_type, slide_idx, total):
    blob = " ".join(s.get("text","") for s in shapes if s.get("has_text")).lower()
    wc = sum(s.get("word_count",0) for s in shapes)
    nc = sum(1 for s in shapes if not s.get("decorative"))
    nk = sum(1 for r in roles if r == "kpi_card")
    ht = any(s.get("has_table") for s in shapes)
    hc = any(s.get("has_chart") for s in shapes)
    di = any(r=="image" and shapes[i]["w_pct"]>40 and shapes[i]["h_pct"]>40 for i,r in enumerate(roles))

    scores = {}
    for intent, ru in INTENT_RULES.items():
        sc = 0.0
        for kw in ru.get("keywords", []):
            if kw.lower() in blob: sc += 0.20
        if ru.get("has_table") and ht: sc += 0.55
        if ru.get("has_chart") and hc: sc += 0.55
        if ru.get("dom_image") and di: sc += 0.45
        if "min_kpi" in ru and nk >= ru["min_kpi"]: sc += 0.45
        if "layouts" in ru and layout_type in ru["layouts"]: sc += 0.25
        if "n_exact" in ru and nc == ru["n_exact"]: sc += 0.20
        if "max_content" in ru and nc <= ru["max_content"]: sc += 0.15
        if "max_words" in ru and wc <= ru["max_words"]: sc += 0.10
        if "min_n" in ru and nc >= ru["min_n"]: sc += 0.10
        if "word_range" in ru and ru["word_range"][0] <= wc <= ru["word_range"][1]: sc += 0.10
        if ru.get("first") and slide_idx == 0: sc += 0.35
        if ru.get("last") and slide_idx == total - 1: sc += 0.25
        scores[intent] = min(sc, 1.0)

    ranked = sorted(scores.items(), key=lambda x: -x[1])
    best, conf = ranked[0]
    runner = ranked[1] if len(ranked) > 1 else (None, 0)
    ambig = (conf - runner[1]) < 0.15
    return {"intent": best, "confidence": round(conf, 3), "ambiguous": ambig,
            "runner_up": runner[0] if ambig else None,
            "needs_llm": conf < 0.65 or ambig,
            "top5": {k: round(v,3) for k,v in ranked[:5]}}


# ═══════════════════════════════════════════════════
# STEP 6 — TOKEN EXTRACTION (KDTREE)
# ═══════════════════════════════════════════════════

def _build_token_tree(theme):
    ct = {k:v for k,v in theme.items() if k.startswith("token.color.")}
    if not ct: return None, {}
    names = list(ct.keys())
    rgb = [[int(ct[n].lstrip("#")[i:i+2],16) for i in (0,2,4)] for n in names]
    return KDTree(np.array(rgb)), {i:n for i,n in enumerate(names)}

def extract_tokens_fast(slide, theme=None):
    colors, fonts = set(), set()
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb: colors.add(str(run.font.color.rgb))
                    except: pass
                    if run.font.name: fonts.add(run.font.name)
        try:
            if hasattr(shape,'fill') and shape.fill.type is not None:
                if shape.fill.fore_color and shape.fill.fore_color.rgb:
                    colors.add(str(shape.fill.fore_color.rgb))
        except: pass

    cmap, unmapped = {}, []
    if theme:
        tree, idx_name = _build_token_tree(theme)
        if tree:
            for h in colors:
                rgb = [int(h[i:i+2],16) for i in (0,2,4)]
                d, i = tree.query(rgb)
                cmap[h] = idx_name[i] if d < 40 else None
                if d >= 40: unmapped.append(h)

    fmap = {}
    if theme:
        hf, bf = theme.get("token.font.heading"), theme.get("token.font.body")
        for f in fonts:
            if hf and f == hf: fmap[f] = "token.font.heading"
            elif bf and f == bf: fmap[f] = "token.font.body"

    return {"raw_colors": list(colors), "raw_fonts": list(fonts),
            "color_map": cmap, "font_map": fmap,
            "unmapped": unmapped, "needs_validation": len(unmapped) > 0}


# ═══════════════════════════════════════════════════
# STEP 7 — BUILD PRIMITIVE (V2 FORMAT)
# ═══════════════════════════════════════════════════

def _quality_score(comps, roles, layout, wc):
    sc = 5.0
    if "title" in roles: sc += 1.0
    else: sc -= 1.5
    for i,a in enumerate(comps):
        for j,b in enumerate(comps):
            if i >= j: continue
            ox = min(a["x_pct"]+a["w_pct"],b["x_pct"]+b["w_pct"]) - max(a["x_pct"],b["x_pct"])
            oy = min(a["y_pct"]+a["h_pct"],b["y_pct"]+b["h_pct"]) - max(a["y_pct"],b["y_pct"])
            if ox > 3 and oy > 3: sc -= 0.5
    if layout in ("grid","columns","rows"): sc += 1.0
    elif layout == "freeform": sc -= 0.5
    if wc > 75: sc -= 1.0
    if len(comps) > 8: sc -= 0.5
    return round(max(0, min(10, sc)), 1)

def build_primitive_v2(classification, shapes, roles, groups, layout,
                       tokens, reading_order, slide_num, source):
    comps = []
    for idx in reading_order:
        s = shapes[idx]
        if s.get("decorative"): continue
        c = {"role": roles[idx], "x_pct": s["x_pct"], "y_pct": s["y_pct"],
             "w_pct": s["w_pct"], "h_pct": s["h_pct"], "rotation": s.get("rotation",0)}
        if s.get("has_text"):
            c["text_preview"] = s.get("text","")[:80]
            c["word_count"] = s.get("word_count",0)
        if s.get("has_table"):
            c["table_size"] = f"{s.get('table_rows',0)}x{s.get('table_cols',0)}"
        comps.append(c)

    tw = sum(c.get("word_count",0) for c in comps)
    q = _quality_score(comps, roles, layout, tw)
    ch = hashlib.sha256(json.dumps(comps, sort_keys=True).encode()).hexdigest()[:12]

    return {
        "schema_version": 2,
        "id": f"prim_{source}_{slide_num}_{ch}",
        "intent": classification["intent"],
        "intent_confidence": classification["confidence"],
        "layout_type": layout,
        "semantic_groups": groups,
        "components": comps,
        "reading_order": list(range(len(comps))),
        "design_tokens": {"colors": tokens.get("color_map",{}),
                          "fonts": tokens.get("font_map",{}),
                          "unmapped": tokens.get("unmapped",[])},
        "density": {"content_shapes": len(comps), "total_words": tw,
                    "decorative_removed": sum(1 for s in shapes if s.get("decorative"))},
        "quality_score": q,
        "status": "draft" if q < 4.0 else "candidate",
        "version": 1,
        "source": {"file": source, "slide": slide_num,
                    "extracted_at": datetime.datetime.now().isoformat(),
                    "extractor": "pisa_v2"},
        "_content_hash": ch,
    }


# ═══════════════════════════════════════════════════
# STEP 8 — DEDUPLICATION (LSH PRE-FILTER)
# ═══════════════════════════════════════════════════

def _pos_sig(comps):
    return frozenset((int(c["x_pct"]//10), int(c["y_pct"]//10)) for c in comps)

def layout_similarity_v2(a, b):
    ca, cb = a["components"], b["components"]
    if abs(len(ca)-len(cb)) > 1: return 0.0
    im = 1.0 if a["intent"] == b["intent"] else 0.0
    lm = 1.0 if a.get("layout_type") == b.get("layout_type") else 0.5
    if not ca or not cb: return 0.0
    sh, lo = (ca,cb) if len(ca)<=len(cb) else (cb,ca)
    td, used = 0, set()
    for x in sh:
        bd, bj = 999, -1
        for j, y in enumerate(lo):
            if j in used: continue
            d = (abs(x["x_pct"]-y["x_pct"])+abs(x["y_pct"]-y["y_pct"])+
                 abs(x["w_pct"]-y["w_pct"])+abs(x["h_pct"]-y["h_pct"]))/4
            d += 0 if x.get("role")==y.get("role") else 5
            if d < bd: bd, bj = d, j
        td += bd
        if bj >= 0: used.add(bj)
    ps = max(0, 1 - (td/len(sh))/40)
    return round(0.15*im + 0.10*lm + 0.75*ps, 3)


# ═══════════════════════════════════════════════════
# STEP 9 — VERSIONED LIBRARY REGISTRATION
# ═══════════════════════════════════════════════════

def register_primitive_v2(prim, lib_path="pisa_library.json", mode="add"):
    if os.path.exists(lib_path):
        with open(lib_path) as f: lib = json.load(f)
    else:
        lib = {"schema_version":2,"primitives":[],"meta":{"version":1,"updated":None}}

    # Same-source versioning
    for ex in lib["primitives"]:
        sa, sb = ex.get("source",{}), prim.get("source",{})
        if sa.get("file")==sb.get("file") and sa.get("slide")==sb.get("slide"):
            prim["version"] = ex.get("version",1)+1
            prim["previous_hash"] = ex.get("_content_hash")
            lib["primitives"].remove(ex)
            lib["primitives"].append(prim)
            lib["meta"]["updated"] = prim["source"]["extracted_at"]
            lib["meta"]["version"] += 1
            with open(lib_path,"w") as f: json.dump(lib,f,indent=2)
            return {"action":"versioned","id":prim["id"],"version":prim["version"]}

    # Dedup
    psig = _pos_sig(prim["components"])
    for ex in lib["primitives"]:
        esig = _pos_sig(ex["components"])
        union = len(psig | esig)
        if union > 0 and len(psig & esig)/union < 0.4: continue
        sim = layout_similarity_v2(prim, ex)
        if sim > DEDUP_THRESHOLD:
            if mode == "add":
                return {"action":"skipped","existing":ex["id"],"similarity":sim}
            elif mode == "replace":
                lib["primitives"].remove(ex); break

    lib["primitives"].append(prim)
    lib["meta"]["updated"] = prim["source"]["extracted_at"]
    lib["meta"]["version"] += 1
    with open(lib_path,"w") as f: json.dump(lib,f,indent=2)
    return {"action":"added","id":prim["id"],"quality":prim["quality_score"],
            "size":len(lib["primitives"])}


# ═══════════════════════════════════════════════════
# STEP 10 — SVG PREVIEW GENERATION (no external deps)
# ═══════════════════════════════════════════════════

def generate_svg_preview(prim, theme):
    """Generate an SVG preview string for a primitive. Pure Python, no external deps.
    Imports the SVG renderer lazily to avoid circular deps."""
    try:
        import sys
        # Add project root to path for import
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        if project_root not in sys.path:
            sys.path.insert(0, project_root)
        from services.svg.primitive_to_svg import render_svg
        return {"svg": render_svg(prim, theme, 960, 540, show_labels=False, show_groups=True)}
    except ImportError:
        return {"svg": None, "note": "SVG renderer not available — install services/svg/primitive_to_svg.py"}


# ═══════════════════════════════════════════════════
# STEP 11 — BATCH ORCHESTRATOR
# ═══════════════════════════════════════════════════

def extract_deck_v2(pptx_path, theme=None, lib_path="pisa_library.json", generate_previews=True):
    prs = Presentation(pptx_path)
    src = os.path.basename(pptx_path).replace(".pptx","")
    total = len(prs.slides)
    results = []
    for i, slide in enumerate(prs.slides):
        shapes = deep_inventory(slide, i)
        content = [s for s in shapes if not s.get("decorative")]
        axes = detect_alignment_axes(content)
        layout = classify_layout(content, axes)
        order = infer_reading_order(content)
        roles = assign_roles(content, order, layout)
        groups = detect_semantic_groups(content, roles, order)
        cls = classify_intent_v2(content, roles, groups, layout, i, total)
        tokens = extract_tokens_fast(slide, theme)
        prim = build_primitive_v2(cls, content, roles, groups, layout, tokens, order, i+1, src)
        reg = register_primitive_v2(prim, lib_path)
        svg = None
        if generate_previews and theme and reg["action"] in ("added","versioned"):
            svg_result = generate_svg_preview(prim, theme)
            svg = svg_result.get("svg")
            # Store SVG in the primitive within the library
            if svg:
                prim["svg"] = svg
        results.append({"slide":i+1,"layout":layout,"intent":cls["intent"],
                        "confidence":cls["confidence"],"ambiguous":cls["ambiguous"],
                        "quality":prim["quality_score"],"components":len(prim["components"]),
                        "decorative_removed":prim["density"]["decorative_removed"],
                        "tokens_ok":not tokens["needs_validation"],
                        "registration":reg,"has_preview":svg is not None})
    return {"source":src,"total":total,"results":results}


# ═══════════════════════════════════════════════════
# CLI ENTRY POINT
# ═══════════════════════════════════════════════════

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="PISA V2 Slide-to-Primitive Extractor")
    parser.add_argument("pptx", help="Path to source PPTX")
    parser.add_argument("--theme", help="Path to theme JSON")
    parser.add_argument("--library", default="pisa_library.json")
    parser.add_argument("--no-previews", action="store_true", help="Skip SVG preview generation")
    args = parser.parse_args()

    theme = None
    if args.theme:
        with open(args.theme) as f: theme = json.load(f)

    result = extract_deck_v2(args.pptx, theme, args.library, not args.no_previews)
    print(json.dumps(result, indent=2))
